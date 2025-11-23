#!/usr/bin/env python3
"""
PowerPoint Generator for Fidelity PlanAlign Studio
Creates peer comparison presentation decks
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from typing import Dict
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import io
from pathlib import Path

from peer_benchmarking import generate_peer_comparison


def generate_powerpoint(client_id: str, output_dir: str = 'output') -> str:
    """
    Generate PowerPoint deck with peer comparison analysis.

    Args:
        client_id: Target client ID
        output_dir: Directory for output files

    Returns:
        Path to generated .pptx file
    """
    # Get peer comparison data
    comparison = generate_peer_comparison(client_id)
    target = comparison['target_client']
    cohort_size = comparison['peer_cohort']['size']

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = target['client_name']
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "401(k) Plan Analysis"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Report type
    report_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
    report_frame = report_box.text_frame
    report_frame.text = "Peer Benchmarking Report"
    report_para = report_frame.paragraphs[0]
    report_para.font.size = Pt(20)
    report_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"Generated: {datetime.now().strftime('%B %Y')}"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Plan Overview & Peer Comparison
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Plan Overview & Peer Comparison"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Plan Details
    details_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4), Inches(2.5))
    details_frame = details_box.text_frame
    details_frame.text = "PLAN DETAILS\n"
    details_frame.paragraphs[0].font.size = Pt(16)
    details_frame.paragraphs[0].font.bold = True

    details_frame.add_paragraph()
    p = details_frame.paragraphs[1]
    p.text = f"Industry: {target['industry']}"
    p.font.size = Pt(14)

    details_frame.add_paragraph()
    p = details_frame.paragraphs[2]
    p.text = f"Employees: {target['employee_count']:,}"
    p.font.size = Pt(14)

    # Match info
    if target.get('match_formula'):
        details_frame.add_paragraph()
        p = details_frame.paragraphs[3]
        p.text = f"Match: {target['match_formula']}"
        p.font.size = Pt(14)

    # Auto-enrollment info
    if target.get('auto_enrollment_enabled'):
        details_frame.add_paragraph()
        p = details_frame.paragraphs[-1]
        ae_rate = target.get('auto_enrollment_rate')
        if ae_rate is not None:
            p.text = f"Auto-Enroll: {ae_rate*100:.0f}%"
        else:
            p.text = "Auto-Enroll: Enabled"
        p.font.size = Pt(14)

    # Peer Cohort Info
    cohort_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(4), Inches(0.8))
    cohort_frame = cohort_box.text_frame
    cohort_frame.text = f"PEER COHORT (n={cohort_size} similar plans)\nSelected based on industry and size"
    cohort_frame.paragraphs[0].font.size = Pt(12)
    cohort_frame.paragraphs[0].font.bold = True

    # Generate chart: Employee Count Comparison
    if 'employee_count' in comparison['numeric_comparisons']:
        employee_comparison = comparison['numeric_comparisons']['employee_count']

        fig, ax = plt.subplots(figsize=(6, 2.5))

        # Peer percentiles
        p25 = employee_comparison['peer_p25']
        p50 = employee_comparison['peer_median']
        p75 = employee_comparison['peer_p75']
        your_value = employee_comparison['your_value']

        # Bar for peer range
        ax.barh([0], [p75 - p25], left=p25, height=0.5, color='lightgray', label='Peer Range (P25-P75)')

        # Median line
        ax.axvline(p50, color='red', linestyle='--', linewidth=2, label='Peer Median')

        # Your value marker
        ax.scatter([your_value], [0], s=200, color='blue', marker='D', zorder=3, label='Your Plan')

        ax.set_yticks([])
        ax.set_xlabel('Employee Count', fontsize=12)
        ax.set_title('Employee Count Comparison', fontsize=14, fontweight='bold')
        ax.legend(loc='upper right', fontsize=9)
        ax.grid(axis='x', alpha=0.3)

        # Save chart to bytes
        chart_stream = io.BytesIO()
        plt.tight_layout()
        plt.savefig(chart_stream, format='png', dpi=150, bbox_inches='tight')
        chart_stream.seek(0)
        plt.close()

        # Add chart to slide
        slide2.shapes.add_picture(chart_stream, Inches(5), Inches(1), width=Inches(4.5))

        # Employee rank text
        rank = employee_comparison['percentile_rank']
        rank_box = slide2.shapes.add_textbox(Inches(5), Inches(3.8), Inches(4.5), Inches(0.8))
        rank_frame = rank_box.text_frame
        rank_frame.text = f"Your Plan: {your_value:,} employees ({rank['percentile']}th percentile)\n{rank['label']}"
        rank_frame.paragraphs[0].font.size = Pt(12)
        rank_frame.paragraphs[0].font.bold = True

    # Feature Adoption
    adoption_box = slide2.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(2))
    adoption_frame = adoption_box.text_frame
    adoption_frame.text = "FEATURE ADOPTION\n"
    adoption_frame.paragraphs[0].font.size = Pt(14)
    adoption_frame.paragraphs[0].font.bold = True

    # Match adoption
    if 'has_match' in comparison['adoption_comparisons']:
        match_adoption = comparison['adoption_comparisons']['has_match']
        adoption_frame.add_paragraph()
        p = adoption_frame.paragraphs[-1]
        match_symbol = "✓" if match_adoption['your_value'] else "✗"
        p.text = f"{match_symbol} Employer Match: {'Offered' if match_adoption['your_value'] else 'Not Offered'} | {match_adoption['peer_adoption_rate']*100:.0f}% of peers"
        p.font.size = Pt(12)

    # Auto-enrollment adoption
    if 'has_auto_enrollment' in comparison['adoption_comparisons']:
        ae_adoption = comparison['adoption_comparisons']['has_auto_enrollment']
        adoption_frame.add_paragraph()
        p = adoption_frame.paragraphs[-1]
        ae_symbol = "✓" if ae_adoption['your_value'] else "✗"
        p.text = f"{ae_symbol} Auto-Enrollment: {'Enabled' if ae_adoption['your_value'] else 'Not Enabled'} | {ae_adoption['peer_adoption_rate']*100:.0f}% of peers"
        p.font.size = Pt(12)

    # SLIDE 3: Key Findings & Recommendations
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Key Findings & Recommendations"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Generate findings and recommendations
    findings = generate_findings(comparison)
    recommendations = generate_recommendations(comparison)

    # Key Findings
    findings_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2.5))
    findings_frame = findings_box.text_frame
    findings_frame.text = "KEY FINDINGS\n"
    findings_frame.paragraphs[0].font.size = Pt(18)
    findings_frame.paragraphs[0].font.bold = True

    for finding in findings:
        findings_frame.add_paragraph()
        p = findings_frame.paragraphs[-1]
        p.text = f"• {finding}"
        p.font.size = Pt(14)
        p.level = 0

    # Recommendations
    rec_box = slide3.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(2.5))
    rec_frame = rec_box.text_frame
    rec_frame.text = "RECOMMENDATIONS\n"
    rec_frame.paragraphs[0].font.size = Pt(18)
    rec_frame.paragraphs[0].font.bold = True

    for rec in recommendations:
        rec_frame.add_paragraph()
        p = rec_frame.paragraphs[-1]
        p.text = f"• {rec}"
        p.font.size = Pt(14)
        p.level = 0

    # Save presentation
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    filename = f"{target['client_name'].replace(' ', '_')}_Peer_Analysis_{datetime.now().strftime('%Y%m%d')}.pptx"
    filepath = Path(output_dir) / filename

    prs.save(str(filepath))

    return str(filepath)


def generate_findings(comparison: Dict) -> list:
    """Generate key findings using simple rules-based logic"""
    findings = []
    target = comparison['target_client']

    # Employee count finding
    if 'employee_count' in comparison['numeric_comparisons']:
        employee_rank = comparison['numeric_comparisons']['employee_count']['percentile_rank']
        findings.append(
            f"Your plan has {target['employee_count']:,} employees "
            f"({employee_rank['label'].lower()}, {employee_rank['percentile']}th percentile) "
            f"compared to similar plans"
        )

    # Match finding
    if target.get('match_formula'):
        if 'match_effective_rate' in comparison['numeric_comparisons']:
            match_rank = comparison['numeric_comparisons']['match_effective_rate']['percentile_rank']
            match_rate = target.get('match_effective_rate')
            if match_rate is not None:
                findings.append(
                    f"Match rate of {match_rate*100:.1f}% is {match_rank['label'].lower()} "
                    f"({match_rank['percentile']}th percentile) among peers offering match"
                )
    else:
        if 'has_match' in comparison['adoption_comparisons']:
            match_adoption = comparison['adoption_comparisons']['has_match']['peer_adoption_rate']
            findings.append(
                f"Your plan does not offer employer match, while {match_adoption*100:.0f}% of peer plans do"
            )

    # Auto-enrollment finding
    if 'has_auto_enrollment' in comparison['adoption_comparisons']:
        ae_adoption = comparison['adoption_comparisons']['has_auto_enrollment']['peer_adoption_rate']
        if target.get('auto_enrollment_enabled'):
            findings.append(
                f"{ae_adoption*100:.0f}% of peer plans have adopted auto-enrollment, "
                f"positioning it as an industry standard"
            )
        else:
            findings.append(
                f"Your plan has not adopted auto-enrollment, though {ae_adoption*100:.0f}% of peers have"
            )

    return findings


def generate_recommendations(comparison: Dict) -> list:
    """Generate recommendations using simple rules-based logic"""
    recommendations = []
    target = comparison['target_client']

    # Match recommendation
    if target.get('match_formula'):
        if 'match_effective_rate' in comparison['numeric_comparisons']:
            match_rank = comparison['numeric_comparisons']['match_effective_rate']['percentile_rank']
            percentile = match_rank.get('percentile')
            if percentile is not None and percentile < 50:
                recommendations.append(
                    "Consider increasing match rate to align with peer median or top quartile"
                )
            elif percentile is not None:
                recommendations.append(
                    "Maintain current match rate which is competitive with peers"
                )
    else:
        if 'has_match' in comparison['adoption_comparisons']:
            match_adoption = comparison['adoption_comparisons']['has_match']['peer_adoption_rate']
            if match_adoption > 0.7:
                recommendations.append(
                    "Consider adding employer match, as it is offered by majority of peer plans"
                )

    # Auto-enrollment recommendation
    if 'has_auto_enrollment' in comparison['adoption_comparisons']:
        ae_adoption = comparison['adoption_comparisons']['has_auto_enrollment']['peer_adoption_rate']
        if not target.get('auto_enrollment_enabled') and ae_adoption > 0.5:
            recommendations.append(
                "Consider implementing auto-enrollment to align with peer practices and improve participation"
            )
        elif target.get('auto_enrollment_enabled'):
            ae_rate = target.get('auto_enrollment_rate')
            if ae_rate is not None:
                recommendations.append(
                    f"Maintain current auto-enrollment design ({ae_rate*100:.0f}% default) "
                    "which aligns with peer practices"
                )
            else:
                recommendations.append(
                    "Maintain current auto-enrollment design which aligns with peer practices"
                )

    # Generic recommendation
    recommendations.append(
        "Continue monitoring peer trends and consider annual benchmarking reviews"
    )

    return recommendations